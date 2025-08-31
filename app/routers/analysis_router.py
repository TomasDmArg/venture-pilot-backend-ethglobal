from fastapi import APIRouter, HTTPException, BackgroundTasks, File, UploadFile, Form
from app.models.schemas import (
    ProjectAnalysisRequest, 
    AnalysisResponse, 
    GitRollScanRequest, 
    GitRollScanResponse
)
from app.services.simple_analysis_service import SimpleAnalysisService
from app.services.viability_agent import ViabilityAgent
from app.services.enhanced_founder_search import EnhancedFounderSearch
from app.services.github_analyzer import GitHubAnalyzer
from app.services.summary_generator import SummaryGenerator
from app.services.google_search_service import GoogleSearchService
from app.services.gitroll_service import GitRollService
from app.services.competitor_agent import CompetitorAgent
from app.services.followup_agent import FollowUpAgent
from app.services.compliance_agent import ComplianceAgent
from typing import Dict, Any, Optional
import asyncio
import base64
import logging
import time
import tiktoken
import io
import json
from openai import OpenAI
from docx import Document as DocxDocument
from pdfminer.high_level import extract_text as extract_pdf_text

# Setup logging
logger = logging.getLogger(__name__)

router = APIRouter(prefix="/analysis", tags=["Project Analysis"])

# Initialize services
analysis_service = SimpleAnalysisService()
viability_agent = ViabilityAgent()
founder_search = EnhancedFounderSearch()
github_analyzer = GitHubAnalyzer()
summary_generator = SummaryGenerator()
gitroll_service = GitRollService()
competitor_agent = CompetitorAgent()
followup_agent = FollowUpAgent()
compliance_agent = ComplianceAgent()

@router.post("/project", response_model=AnalysisResponse)
async def analyze_project(request: ProjectAnalysisRequest):
    """
    Analyze a project deck and find founders information
    """
    try:
        # Validate input
        if not request.deck_file:
            raise HTTPException(status_code=400, detail="Deck file is required")
        
        # Perform analysis
        project_summary = await analysis_service.analyze_project(
            deck_file=request.deck_file,
            project_name=request.project_name
        )
        
        return AnalysisResponse(
            success=True,
            message="Project analysis completed successfully",
            data=project_summary
        )
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")

@router.post("/project/upload")
async def analyze_project_upload(
    file: UploadFile = File(...),
    project_name: Optional[str] = Form(None)
):
    """
    Analyze a project deck uploaded as a file (FormData) - Complete Report
    """
    start_time = time.time()
    try:
        logger.info(f"Starting complete analysis for file: {file.filename}")
        
        if not file:
            logger.error("No file provided")
            return {
                "status": "error",
                "message": "No file provided",
                "analysis_completed": False
            }
        
        # Detect file type and extract content
        filename = file.filename or ""
        ext = filename.lower().split('.')[-1]
        logger.info(f"Processing file: {filename} (type: {ext})")
        
        content = await file.read()
        deck_content = None
        
        if ext == "pdf":
            # PDF extraction
            logger.info("Extracting text from PDF...")
            deck_content = extract_pdf_text(io.BytesIO(content))
        elif ext == "pptx":
            # PPTX extraction
            logger.info("Extracting text from PPTX...")
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:  # type: ignore
                        try:
                            slides_text.append(shape.text_frame.text)  # type: ignore
                        except Exception:
                            pass
            deck_content = "\n".join(slides_text)
        elif ext == "docx":
            # DOCX extraction
            logger.info("Extracting text from DOCX...")
            doc = DocxDocument(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs]
            deck_content = "\n".join(paragraphs)
        elif ext in ["txt", "md"]:
            # TXT or Markdown
            logger.info("Extracting text from TXT/MD...")
            try:
                deck_content = content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    deck_content = content.decode('latin-1')
                except UnicodeDecodeError:
                    deck_content = content.decode('cp1252', errors='ignore')
        else:
            logger.error(f"Unsupported file type: {ext}")
            return {
                "status": "error",
                "message": "Unsupported file type. Use PDF, PPTX, DOCX, TXT or MD.",
                "analysis_completed": False
            }
        
        if not deck_content or not deck_content.strip():
            logger.error("Could not extract text from file")
            return {
                "status": "error",
                "message": "Could not extract text from file.",
                "analysis_completed": False
            }
        
        logger.info(f"Text extraction completed. Content length: {len(deck_content)} characters")
        
        # Step 1: Basic project analysis
        logger.info("Step 1: Starting basic project analysis...")
        step1_start = time.time()
        project_summary = await analysis_service.analyze_project_content(
            deck_content=deck_content,
            project_name=project_name or file.filename
        )
        step1_time = time.time() - step1_start
        logger.info(f"Step 1 completed in {step1_time:.1f}s - Project: {project_summary.project_name}")
        
        # Step 2: Extract project info for other analyses
        project_info = {
            "project_name": project_summary.project_name,
            "description": project_summary.description,
            "problem_statement": project_summary.problem_statement,
            "solution": project_summary.solution,
            "target_market": project_summary.target_market,
            "business_model": project_summary.business_model,
            "team_info": project_summary.team_info
        }
        
        # Step 3: Viability assessment
        logger.info("Step 2: Starting viability assessment...")
        step2_start = time.time()
        viability_assessment = await viability_agent.assess_viability(project_info)
        step2_time = time.time() - step2_start
        logger.info(f"Step 2 completed in {step2_time:.1f}s - Viability Score: {viability_assessment.get('score', 'N/A')}")
        
        # Step 4: Founder search
        logger.info("Step 3: Starting founder search and analysis...")
        step3_start = time.time()
        founders = await founder_search.extract_and_search_founders(deck_content, project_summary.project_name)
        step3_time = time.time() - step3_start
        logger.info(f"Step 3 completed in {step3_time:.1f}s - Founders: {len(founders)}")
        
        # Step 5: GitHub analysis
        logger.info("Step 4: Starting GitHub analysis...")
        step4_start = time.time()
        github_repos = await github_analyzer.analyze_github_repos(deck_content, founders)
        step4_time = time.time() - step4_start
        logger.info(f"Step 4 completed in {step4_time:.1f}s - GitHub repos: {len(github_repos)}")
        
        # Step 6: Generate simple summary
        logger.info("Step 5: Generating summary...")
        step5_start = time.time()
        simple_summary = await summary_generator.generate_simple_summary(
            project_info, 
            viability_assessment.get("score", 5)
        )
        step5_time = time.time() - step5_start
        logger.info(f"Step 5 completed in {step5_time:.1f}s")
        
        # Step 7: Competitor analysis
        logger.info("Step 6: Starting competitor analysis...")
        step6_start = time.time()
        competitor_analysis = await competitor_agent.analyze_competitors(project_info)
        step6_time = time.time() - step6_start
        logger.info(f"Step 6 completed in {step6_time:.1f}s - Competitors found: {len(competitor_analysis.get('competitors', []))}")
        
        # Step 8: Compliance analysis
        logger.info("Step 7: Starting compliance analysis...")
        step7_start = time.time()
        compliance_analysis = await compliance_agent.analyze_compliance_risks(project_info, founders)
        step7_time = time.time() - step7_start
        logger.info(f"Step 7 completed in {step7_time:.1f}s - Compliance score: {compliance_analysis.get('compliance_score', 'N/A')}")
        
        # Step 9: Generate follow-up questions
        logger.info("Step 8: Generating follow-up questions...")
        step8_start = time.time()
        followup_questions = await followup_agent.generate_followup_questions(
            project_info, 
            founders, 
            viability_assessment, 
            competitor_analysis
        )
        step8_time = time.time() - step8_start
        logger.info(f"Step 8 completed in {step8_time:.1f}s - Questions generated: {followup_questions.get('total_questions', 0)}")
        
        # Create complete response
        total_time = time.time() - start_time
        logger.info(f"Complete analysis finished in {total_time:.1f}s")
        
        complete_response = {
            "status": "success",
            "project_name": project_summary.project_name,
            "summary": simple_summary,
            "viability_score": viability_assessment.get("score", 5),
            "viability_explanation": viability_assessment.get("explanation", ""),
            "viability_breakdown": {
                "team_score": viability_assessment.get("team_score", 5),
                "market_score": viability_assessment.get("market_score", 5),
                "product_score": viability_assessment.get("product_score", 5),
                "business_model_score": viability_assessment.get("business_model_score", 5),
                "execution_score": viability_assessment.get("execution_score", 5)
            },
            "risk_factors": viability_assessment.get("risk_factors", []),
            "strengths": viability_assessment.get("strengths", []),
            "penalties_applied": viability_assessment.get("penalties_applied", []),
            "bonuses_applied": viability_assessment.get("bonuses_applied", []),
            "critical_concerns": viability_assessment.get("critical_concerns", []),
            "recommendation": viability_assessment.get("recommendation", "More research needed"),
            "founders": founders,
            "github_analysis": github_repos,
            "competitor_analysis": {
                "competitors": competitor_analysis.get("competitors", []),
                "market_saturation": competitor_analysis.get("market_saturation", "unknown"),
                "competitive_advantage": competitor_analysis.get("competitive_advantage", ""),
                "threat_level": competitor_analysis.get("threat_level", "medium"),
                "key_differentiators": competitor_analysis.get("competitive_analysis", {}).get("key_differentiators", []),
                "market_gaps": competitor_analysis.get("competitive_analysis", {}).get("market_gaps", [])
            },
            "compliance_analysis": {
                "compliance_score": compliance_analysis.get("compliance_score", 5),
                "risk_level": compliance_analysis.get("risk_level", "medium"),
                "compliance_risks": compliance_analysis.get("compliance_risks", []),
                "regulatory_requirements": compliance_analysis.get("regulatory_requirements", []),
                "legal_risks": compliance_analysis.get("legal_risks", []),
                "data_privacy_concerns": compliance_analysis.get("data_privacy_concerns", []),
                "required_licenses": compliance_analysis.get("required_licenses", [])
            },
            "followup_questions": {
                "questions": followup_questions.get("followup_questions", []),
                "total_questions": followup_questions.get("total_questions", 0),
                "categories": followup_questions.get("categories", {}),
                "priority_questions": followup_questions.get("priority_questions", [])
            },
            "detailed_analysis": {
                "description": project_summary.description,
                "problem": project_summary.problem_statement,
                "solution": project_summary.solution,
                "target_market": project_summary.target_market,
                "business_model": project_summary.business_model
            },
            "analysis_completed": True,
            "message": f"Complete analysis performed for: {project_summary.project_name}",
            "file_processed": file.filename,
            "processing_time_seconds": round(total_time, 1),
            "step_timing": {
                "step1_project_analysis": round(step1_time, 1),
                "step2_viability_assessment": round(step2_time, 1),
                "step3_founder_search": round(step3_time, 1),
                "step4_github_analysis": round(step4_time, 1),
                "step5_summary_generation": round(step5_time, 1),
                "step6_competitor_analysis": round(step6_time, 1),
                "step7_compliance_analysis": round(step7_time, 1),
                "step8_followup_questions": round(step8_time, 1)
            }
        }
        
        return complete_response
        
    except Exception as e:
        total_time = time.time() - start_time
        logger.error(f"Error during analysis after {total_time:.1f}s: {str(e)}")
        return {
            "status": "error", 
            "message": f"Error during analysis: {str(e)}",
            "analysis_completed": False,
            "processing_time_seconds": round(total_time, 1)
        }

@router.post("/project/simple")
async def analyze_project_simple(request: ProjectAnalysisRequest):
    """
    Analyze a project deck and return human-readable response
    """
    try:
        # Validate input
        if not request.deck_file:
            raise HTTPException(status_code=400, detail="Deck file is required")
        
        # Perform analysis
        project_summary = await analysis_service.analyze_project(
            deck_file=request.deck_file,
            project_name=request.project_name
        )
        
        # Create human-readable response
        human_response = {
            "status": "success",
            "project_name": project_summary.project_name,
            "summary": {
                "description": project_summary.description,
                "problem": project_summary.problem_statement,
                "solution": project_summary.solution,
                "target_market": project_summary.target_market,
                "business_model": project_summary.business_model
            },
            "analysis_completed": True,
            "message": f"Analysis completed for project: {project_summary.project_name}"
        }
        
        return human_response
        
    except ValueError as e:
        return {
            "status": "error",
            "message": f"Validation error: {str(e)}",
            "analysis_completed": False
        }
    except Exception as e:
        return {
            "status": "error", 
            "message": f"Error during analysis: {str(e)}",
            "analysis_completed": False
        }

@router.post("/gitroll/scan", response_model=GitRollScanResponse)
async def initiate_gitroll_scan(request: GitRollScanRequest):
    """
    Initiate a GitRoll scan for a GitHub user
    """
    try:
        if not request.user:
            raise HTTPException(status_code=400, detail="GitHub username is required")
        
        result = await gitroll_service.initiate_scan(request.user)
        
        return GitRollScanResponse(
            success=result.success,
            scan_id=result.scan_id,
            profile_url=result.profile_url,
            message=result.message
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"GitRoll scan failed: {str(e)}")

@router.get("/gitroll/status/{scan_id}")
async def check_gitroll_status(scan_id: str):
    """
    Check the status of a GitRoll scan
    """
    try:
        status = await gitroll_service.check_scan_status(scan_id)
        
        if status.get("success"):
            return {
                "success": True,
                "scan_id": scan_id,
                "status": status.get("status", "unknown"),
                "score": status.get("score"),
                "og_image_score": status.get("og_image_score"),
                "profile_url": status.get("profile_url")
            }
        else:
            return {
                "success": False,
                "scan_id": scan_id,
                "message": status.get("message", "Unknown error")
            }
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error checking scan status: {str(e)}")

@router.get("/health")
async def health_check():
    """
    Health check endpoint
    """
    return {"status": "healthy", "service": "analysis"}

@router.get("/")
async def root():
    """
    Root endpoint for analysis router
    """
    return {
        "message": "Project Analysis API",
        "endpoints": {
            "analyze_project": "POST /analysis/project",
            "analyze_project_upload": "POST /analysis/project/upload",
            "analyze_project_simple": "POST /analysis/project/simple",
            "followup_questions": "POST /analysis/followup-questions",
            "gitroll_scan": "POST /analysis/gitroll/scan",
            "gitroll_status": "GET /analysis/gitroll/status/{scan_id}",
            "health": "GET /analysis/health"
        },
        "agents": {
            "project_analysis": "Basic project analysis and information extraction",
            "viability_assessment": "Strict VC-style viability scoring with penalties/bonuses",
            "founder_search": "Comprehensive founder analysis with scoring and GitRoll integration",
            "github_analyzer": "GitHub repository and code analysis",
            "competitor_agent": "Competitive landscape analysis and market saturation",
            "compliance_agent": "Legal, regulatory, and compliance risk assessment",
            "followup_agent": "Top-tier follow-up questions for due diligence",
            "summary_generator": "One-line project summary generation"
        }
    }

@router.post("/followup-questions")
async def generate_followup_questions_endpoint(
    file: UploadFile = File(...),
    project_name: Optional[str] = Form(None)
):
    """
    Generate follow-up questions for a project deck
    """
    start_time = time.time()
    try:
        logger.info(f"Generating follow-up questions for file: {file.filename}")
        
        if not file:
            return {
                "status": "error",
                "message": "No file provided",
                "questions_generated": False
            }
        
        # Extract text from file (simplified version)
        filename = file.filename or ""
        ext = filename.lower().split('.')[-1]
        content = await file.read()
        deck_content = None
        
        if ext in ["txt", "md"]:
            try:
                deck_content = content.decode('utf-8')
            except UnicodeDecodeError:
                deck_content = content.decode('latin-1')
        else:
            return {
                "status": "error",
                "message": "Currently only supports TXT and MD files for follow-up questions",
                "questions_generated": False
            }
        
        if not deck_content or not deck_content.strip():
            return {
                "status": "error",
                "message": "Could not extract text from file",
                "questions_generated": False
            }
        
        # Basic project analysis
        project_summary = await analysis_service.analyze_project_content(
            deck_content=deck_content,
            project_name=project_name or file.filename
        )
        
        project_info = {
            "project_name": project_summary.project_name,
            "description": project_summary.description,
            "problem_statement": project_summary.problem_statement,
            "solution": project_summary.solution,
            "target_market": project_summary.target_market,
            "business_model": project_summary.business_model
        }
        
        # Extract founders
        founders = await founder_search.extract_and_search_founders(deck_content, project_summary.project_name)
        
        # Basic viability assessment
        viability_assessment = await viability_agent.assess_viability(project_info)
        
        # Basic competitor analysis
        competitor_analysis = await competitor_agent.analyze_competitors(project_info)
        
        # Generate follow-up questions
        followup_questions = await followup_agent.generate_followup_questions(
            project_info, 
            founders, 
            viability_assessment, 
            competitor_analysis
        )
        
        total_time = time.time() - start_time
        logger.info(f"Follow-up questions generated in {total_time:.1f}s")
        
        return {
            "status": "success",
            "project_name": project_summary.project_name,
            "followup_questions": followup_questions.get("followup_questions", []),
            "total_questions": followup_questions.get("total_questions", 0),
            "categories": followup_questions.get("categories", {}),
            "priority_questions": followup_questions.get("priority_questions", []),
            "questions_generated": True,
            "processing_time_seconds": round(total_time, 1)
        }
        
    except Exception as e:
        total_time = time.time() - start_time
        logger.error(f"Error generating follow-up questions after {total_time:.1f}s: {str(e)}")
        return {
            "status": "error",
            "message": f"Error generating follow-up questions: {str(e)}",
            "questions_generated": False,
            "processing_time_seconds": round(total_time, 1)
        }

@router.post("/document/analyze")
async def analyze_document(
    file: UploadFile = File(...)
):
    """
    Analyze any VC document (Term Sheet, SAFE, SAFT, SPA, Shareholders' Agreement, Cap Table, Due Diligence, KYC, etc.).
    Returns risk summary and score.
    """
    start_time = time.time()
    # 1. Extract text from file
    content = await file.read()
    if not file.filename:
        return {"error": "File must have a filename."}
    ext = file.filename.lower().split('.')[-1]
    if ext == "pdf":
        text = extract_pdf_text(io.BytesIO(content))
    elif ext == "docx":
        doc = DocxDocument(io.BytesIO(content))
        text = "\n".join([p.text for p in doc.paragraphs])
    elif ext in ["txt", "md"]:
        try:
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1')
    else:
        return {"error": "Unsupported file type. Use PDF, DOCX, TXT, or MD."}
    if not text or not text.strip():
        return {"error": "Could not extract text from file."}

    # 1.5. Detect document type using OpenAI
    client = OpenAI()
    detect_prompt = f"You are an expert legal analyst. Given the following document text, classify the document type as one of: Term Sheet, SAFE, SAFT, SPA, Shareholders' Agreement, Cap Table, Due Diligence, KYC, or Other. Respond ONLY with the type string, nothing else.\n\n---\n\n{text[:3000]}\n\n---\n\nType:"
    resp = client.chat.completions.create(
        model="gpt-5-nano-2025-08-07",
        messages=[
            {"role": "system", "content": "You are an expert legal analyst."},
            {"role": "user", "content": detect_prompt}
        ],
        temperature=0.1
    )
    doc_type_raw = resp.choices[0].message.content
    if doc_type_raw:
        doc_type = doc_type_raw.strip().split("\n")[0]
        if not doc_type:
            doc_type = "Other"
    else:
        doc_type = "Other"

    # 2. Chunking (by 3000 tokens)
    enc = tiktoken.encoding_for_model("gpt-3.5-turbo")
    tokens = enc.encode(text)
    max_tokens = 3000
    chunks = []
    for i in range(0, len(tokens), max_tokens):
        chunk_tokens = tokens[i:i+max_tokens]
        chunk_text = enc.decode(chunk_tokens)
        chunks.append(chunk_text)

    # 3. Clause extraction (OpenAI call)
    clause_chunks = []
    for chunk in chunks:
        prompt = f"""Extract only the shortest, most risk-relevant clauses from the following {doc_type} document chunk. Ignore boilerplate, generic, or non-risky content. Return a JSON array of concise clause strings, each focused on a specific risk or obligation.\n\n---\n\n{chunk}\n\n---\n\nJSON array of clauses:"""
        resp = client.chat.completions.create(
            model="gpt-5-nano-2025-08-07",
            messages=[
                {"role": "system", "content": "You are an expert VC legal analyst. Only extract clauses that present a real risk or obligation. Be extremely concise."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1
        )
        try:
            content_str = resp.choices[0].message.content or ''
            clauses = json.loads(content_str)
        except Exception:
            import re
            match = re.search(r'\[.*\]', content_str, re.DOTALL)
            if match:
                group = match.group(0)
                if group and isinstance(group, str):
                    clauses = json.loads(group)
                else:
                    clauses = []
            else:
                clauses = []
        clause_chunks.extend(clauses)

    # 4. Risk evaluation (OpenAI call per clause)
    risk_objs = []
    for clause in clause_chunks:
        prompt = f"""For the following clause from a {doc_type} document, respond ONLY if it presents a real risk. Label the risk as 'low', 'medium', or 'critical'. Respond in this JSON format:\n{{\n  \"clause\": \"<1-line risk clause>\",\n  \"risk_level\": \"low\"|\"medium\"|\"critical\"\n}}\nIf there is no real risk, do not return anything.\n\nClause: {clause}\n\nJSON:"""
        resp = client.chat.completions.create(
            model="gpt-5-nano-2025-08-07",
            messages=[
                {"role": "system", "content": "You are an expert VC legal risk analyst. Only respond if there is a real risk. Be extremely concise."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1
        )
        try:
            content_str = resp.choices[0].message.content or ''
            if not content_str.strip():
                continue
            risk_obj = json.loads(content_str)
        except Exception:
            import re
            match = re.search(r'\{.*\}', content_str, re.DOTALL)
            if match:
                group = match.group(0)
                if group and isinstance(group, str):
                    risk_obj = json.loads(group)
                else:
                    continue
            else:
                continue
        risk_objs.append(risk_obj)

    # 5. Score computation (OpenAI call)
    score_prompt = f"""Given the following array of clause risk objects from a {doc_type} document, compute an overall risk score from 0 (very risky) to 100 (very safe).\n\n- Give a score close to 100 ONLY if the document is extremely safe, with no critical risks and very few or no medium risks.\n- If there are any critical risks, or several medium risks, the score should be much lower.\n- Be strict and harsh: do NOT be generous. Even a few real risks should significantly lower the score.\n- Respond in this JSON format: {{ \"score\": <number 0-100> }}\n\nArray of risks:\n{json.dumps(risk_objs, ensure_ascii=False)}\n\nJSON:"""
    resp = client.chat.completions.create(
        model="gpt-5-nano-2025-08-07",
        messages=[
            {"role": "system", "content": "You are an expert VC risk scoring analyst. Be objective and rigorous."},
            {"role": "user", "content": score_prompt}
        ],
        temperature=0.1
    )
    try:
        content_str = resp.choices[0].message.content or ''
        score_obj = json.loads(content_str)
    except Exception:
        import re
        match = re.search(r'\{.*\}', content_str, re.DOTALL)
        if match:
            group = match.group(0)
            if group and isinstance(group, str):
                score_obj = json.loads(group)
            else:
                score_obj = {"score": 50}
        else:
            score_obj = {"score": 50}

    # 6. Build risk summary
    risk_summary = {"critical": [], "medium": [], "low": []}
    for r in risk_objs:
        lvl = r.get("risk_level", "medium")
        if not lvl:
            lvl = "medium"
        lvl = str(lvl).lower()
        if lvl not in risk_summary:
            lvl = "medium"
        risk_summary[lvl].append(r.get("clause", ""))

    # 7. Final response
    return {
        "document_type": doc_type,
        "score": score_obj.get("score", 50),
        "risk_summary": risk_summary
    } 